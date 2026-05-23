# -*- coding: utf-8 -*-
"""
Met à jour le modèle exports/Daily_Market_Report.pptx avec les données de
cotation (exports/cotation.json, issu de /marches/actions) et écrit un fichier
daté exports/Daily_Market_Report_<date>.pptx — le modèle d'origine est conservé.

Stratégie : on modifie le texte des formes nommées (KpiValue1, DatePill...) et
des cellules de tableaux EXISTANTS, en préservant la mise en forme des runs
(on ne réécrit que .text du premier run). Aucune recréation de slide.

Champs absents de l'Excel (valeur transigée FCFA, obligations, variation J/J des
volumes) → lignes retirées de la table INDICES & VOLUMES (choix utilisateur).

Usage : python scripts/update-daily-report.py
"""
import json
import sys
from pathlib import Path
from datetime import date

from pptx import Presentation
from pptx.dml.color import RGBColor

ROOT = Path(__file__).resolve().parent.parent
EXPORTS = ROOT / "exports"
TEMPLATE = EXPORTS / "Daily_Market_Report.pptx"
DATA = EXPORTS / "cotation.json"

data = json.loads(DATA.read_text(encoding="utf-8"))
gen_date = (data.get("session", {}).get("fetchedAt") or date.today().isoformat())[:10]
y, m, d = gen_date.split("-")
date_fr = f"{d}/{m}/{y}"
OUT = EXPORTS / f"Daily_Market_Report_{gen_date}.pptx"

GREEN = RGBColor(0x16, 0xA3, 0x4A)
RED = RGBColor(0xDC, 0x26, 0x26)
GRAY = RGBColor(0x64, 0x74, 0x8B)

# ── Formatage fr-FR ──────────────────────────────────────────────────────────
def fr_int(n):
    return f"{round(n):,.0f}".replace(",", " ")

def fr_dec(n, dec=2):
    return f"{n:,.{dec}f}".replace(",", " ").replace(".", ",")

def signed_pct(v, dec=2):
    sign = "+" if v >= 0 else "-"
    return f"{sign}{abs(v):.{dec}f}".replace(".", ",") + "%"

def pct_color(v):
    return GRAY if abs(v) < 1e-9 else (GREEN if v > 0 else RED)

def kpi_var(v):
    arrow = "▲" if v >= 0 else "▼"  # ▲ / ▼
    return f"{arrow} {signed_pct(v)}"

# ── Indices par code ─────────────────────────────────────────────────────────
idx_by_code = {i["code"]: i for i in data["indices"]}

def idx(code):
    return idx_by_code.get(code)

# ── Helpers texte (préserve la mise en forme du 1er run) ─────────────────────
def set_tf(tf, text, color=None):
    p = tf.paragraphs[0]
    runs = p.runs
    if not runs:
        run = p.add_run()
    else:
        run = runs[0]
        for r in runs[1:]:
            r._r.getparent().remove(r._r)
    run.text = text
    if color is not None:
        run.font.color.rgb = color
    # supprime les paragraphes supplémentaires
    for extra in tf.paragraphs[1:]:
        extra._p.getparent().remove(extra._p)

def set_shape(shape, text, color=None):
    set_tf(shape.text_frame, text, color)

def set_cell(cell, text, color=None):
    set_tf(cell.text_frame, text, color)

# ── Index des formes / tables par nom ────────────────────────────────────────
prs = Presentation(str(TEMPLATE))
shapes_by_name = {}
tables_by_name = {}
for slide in prs.slides:
    for sh in slide.shapes:
        shapes_by_name.setdefault(sh.name, sh)
        if sh.has_table:
            tables_by_name.setdefault(sh.name, sh.table)

def S(name):
    return shapes_by_name[name]

# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 1
# ════════════════════════════════════════════════════════════════════════════
# Dates
for nm in ("DatePill",):
    pass
# Toutes les DatePill (slide1 et slide2) → date du jour
for slide in prs.slides:
    for sh in slide.shapes:
        if sh.name == "DatePill":
            set_shape(sh, date_fr)

# Cartes KPI (COMPOSITE / 30 / PRESTIGE)
kpi_map = [("KpiValue1", "KpiVar1", "BRVMC"),
           ("KpiValue2", "KpiVar2", "BRVM30"),
           ("KpiValue3", "KpiVar3", "BRVMPR")]
for val_name, var_name, code in kpi_map:
    i = idx(code)
    if not i:
        print(f"  ! indice {code} introuvable", file=sys.stderr)
        continue
    set_shape(S(val_name), f"{fr_dec(i['value'])} pts")
    set_shape(S(var_name), kpi_var(i["variationPct"]), color=pct_color(i["variationPct"]))

# Table INDICES & VOLUMES (Tableau 17)
t = tables_by_name["Tableau 17"]
# r1..r4 : indices
row_codes = ["BRVMC", "BRVM30", "BRVMPR", "BRVMPA"]
for r, code in enumerate(row_codes, start=1):
    i = idx(code)
    if not i:
        continue
    set_cell(t.cell(r, 1), f"{fr_dec(i['value'])} pts")
    set_cell(t.cell(r, 2), signed_pct(i["variationPct"]), color=pct_color(i["variationPct"]))
# r5 : volume actions (variation J/J indisponible → tiret), corrige le libellé
set_cell(t.cell(5, 0), "Volume transigé actions")
set_cell(t.cell(5, 1), f"{fr_int(data['marketStats']['totalVolume'])} titres")
set_cell(t.cell(5, 2), "—")  # —
# Retire les lignes non couvertes par l'Excel : r6..r9 (de la fin vers le début)
for r in (9, 8, 7, 6):
    tr = t.rows[r]._tr
    tr.getparent().remove(tr)

# TOP 5 hausses (Tableau 31) / FLOP 5 baisses (Tableau 34)
# Noms courts : retire le suffixe pays, préserve les acronymes, abrège les
# raisons sociales trop longues. Override pour quelques grands comptes.
_COUNTRY = {"CI", "SN", "TG", "BN", "BF", "ML", "NE", "SL", "GW"}
_SMALL = {"de", "du", "des", "la", "le", "les", "et", "of", "pour", "l", "d"}
# Acronymes BRVM à conserver en majuscules (source 100% capitalisée → impossible
# de les détecter automatiquement, d'où la liste blanche).
_ACRONYMS = {"CIE", "SGB", "BOA", "SIB", "NSIA", "SOGB", "BICI", "BIIC",
             "CFAO", "NEI", "CEDA", "SMB", "PALM"}
_NAME_OVERRIDE = {
    "SIBC": "SIB", "BICB": "BIIC Bénin", "CBIBF": "Coris Bank Intl",
    "ETIT": "Ecobank TI", "STBC": "SITAB", "UNLC": "Unilever",
}

def short_name(code, name, maxlen=22):
    if code in _NAME_OVERRIDE:
        return _NAME_OVERRIDE[code]
    toks = name.split()
    while toks and toks[-1].upper() in _COUNTRY:
        toks.pop()
    out = []
    for i, t in enumerate(toks):
        if t.upper() in _ACRONYMS:
            out.append(t.upper())
        elif t.lower() in _SMALL and i > 0:
            out.append(t.lower())
        else:
            out.append(t.capitalize())
    s = " ".join(out)
    return s if len(s) <= maxlen else s[: maxlen - 1].rstrip() + "…"

def fill_movers(table_name, rows, color):
    tbl = tables_by_name[table_name]
    for r in range(1, 6):
        a = rows[r - 1] if r - 1 < len(rows) else None
        if a is None:
            for c in range(4):
                set_cell(tbl.cell(r, c), "")
            continue
        set_cell(tbl.cell(r, 0), a["code"])
        set_cell(tbl.cell(r, 1), short_name(a["code"], a["name"]))
        set_cell(tbl.cell(r, 2), fr_int(a["price"]))
        set_cell(tbl.cell(r, 3), signed_pct(a["changePercent"]), color=color)

fill_movers("Tableau 31", data["topGainers"], GREEN)
fill_movers("Tableau 34", data["topLosers"], RED)

# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 2 — TOP 15 par capitalisation (Tableau 6)
# ════════════════════════════════════════════════════════════════════════════
SECTOR_ABBR = {
    "Télécommunications": "Télécom",
    "Services financiers": "Finance",
    "Consommation de base": "Conso base",
    "Consommation discrétionnaire": "Conso discr.",
    "Industriels": "Industrie",
    "Services aux collectivités": "Services pub.",
    "Services publics": "Services pub.",
    "Énergie": "Énergie",
    "Energie": "Énergie",
    "Matériaux": "Matériaux",
    "Santé": "Santé",
    "Agro-industrie": "Agro",
    "Immobilier": "Immobilier",
}

def abbr_sector(s):
    if s in SECTOR_ABBR:
        return SECTOR_ABBR[s]
    return s if len(s) <= 12 else s[:11] + "."

top15 = sorted(data["actions"], key=lambda a: a["capitalization"], reverse=True)[:15]
tbl = tables_by_name["Tableau 6"]
for r, a in enumerate(top15, start=1):
    set_cell(tbl.cell(r, 0), f"{a['code']} - {short_name(a['code'], a['name'])}")
    set_cell(tbl.cell(r, 1), abbr_sector(a["sector"]))
    set_cell(tbl.cell(r, 2), fr_int(a["capitalization"] / 1e9))
    set_cell(tbl.cell(r, 3), fr_int(a["price"]))
    set_cell(tbl.cell(r, 4), signed_pct(a["changePercent"]), color=pct_color(a["changePercent"]))
    ytd = a.get("ytdPct")
    if ytd is None:
        set_cell(tbl.cell(r, 5), "n.d.", color=GRAY)
    else:
        set_cell(tbl.cell(r, 5), signed_pct(ytd), color=pct_color(ytd))
    set_cell(tbl.cell(r, 6), fr_int(a["volume"]))
    set_cell(tbl.cell(r, 7), f"{fr_dec(a['per'], 1)}x" if a.get("hasPer") else "n.d.")
    set_cell(tbl.cell(r, 8), f"{fr_dec(a['yieldPct'], 1)}%" if a.get("hasYield") else "n.d.")

# Section header slide 2 : nb de valeurs
# (laisse le titre tel quel — "TOP 15")

# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 3 — millésime du copyright
# ════════════════════════════════════════════════════════════════════════════
if "Disclaimer" in shapes_by_name:
    tf = shapes_by_name["Disclaimer"].text_frame
    for p in tf.paragraphs:
        for run in p.runs:
            if "2024" in run.text:
                run.text = run.text.replace("2024", y)

prs.save(str(OUT))
print(f"✓ Rapport mis à jour : {OUT}")
print(f"  Séance : {data['session'].get('sessionLabel')} | indices {len(data['indices'])} | actions {len(data['actions'])}")
